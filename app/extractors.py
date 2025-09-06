from pptx import Presentation
from PyPDF2 import PdfReader
import tempfile


def extract_text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    out = []
    for p in reader.pages:
        out.append(p.extract_text() or "")
    return "\n\n".join(out)


def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return "\n\n".join(slides)
